import os
import zipfile
import pandas as pd
from PyPDF2 import PdfReader
import docx
import json
import csv
from PIL import Image
import speech_recognition as sr
import pytesseract
import tempfile
from pptx import Presentation
from Bio import PDB
from pydub import AudioSegment

from dotenv import load_dotenv
import psycopg2
from psycopg2 import sql

# Define the extraction functions for all supported file types
def extract_text_from_pdf(file_path):
    try:
        with open(file_path, 'rb') as file:
            reader = PdfReader(file)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
        return text
    except Exception as e:
        return f"Error reading PDF file: {e}"

def extract_text_from_docx(file_path):
    try:
        doc = docx.Document(file_path)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        return text
    except Exception as e:
        return f"Error reading DOCX file: {e}"

def extract_text_from_excel(file_path):
    try:
        df = pd.read_excel(file_path, engine='openpyxl')
        text = df.to_string(index=False)
        return text
    except Exception as e:
        return f"Error reading Excel file: {e}. Make sure 'openpyxl' is installed."

def extract_text_from_csv(file_path):
    try:
        with open(file_path, 'r', newline='', encoding='utf-8') as file:
            reader = csv.reader(file)
            text = "\n".join([", ".join(row) for row in reader])
        return text
    except Exception as e:
        return f"Error reading CSV file: {e}"

def extract_text_from_json(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            data = json.load(file)
            text = json.dumps(data, indent=4)
        return text
    except Exception as e:
        return f"Error reading JSON file: {e}"
    
def extract_text_from_jsonld(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            data = json.load(file)
            text = json.dumps(data, indent=4)  # Format the JSON-LD data as a pretty-printed string
        return text
    except Exception as e:
        return f"Error reading JSON-LD file: {e}"    
    
# Function to extract text from .jsonl files
def extract_text_from_jsonl(file_path):
    try:
        text = ""
        with open(file_path, 'r', encoding='utf-8') as file:
            for line in file:
                try:
                    json_object = json.loads(line)
                    text += json.dumps(json_object, indent=4) + "\n"
                except json.JSONDecodeError as e:
                    text += f"Invalid JSON in line: {line}\nError: {e}\n"
        return text
    except Exception as e:
        return f"Error reading JSONL file: {e}"

def extract_text_from_txt(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            text = file.read()
        return text
    except Exception as e:
        return f"Error reading TXT file: {e}"

def extract_text_from_image(file_path):
    try:
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img)
        return text
    except Exception as e:
        return f"Error reading image file: {e}"


# Convert any audio file format to WAV using pydub
def convert_to_wav(audio_path):
    audio_extension = os.path.splitext(audio_path)[1].lower()
    if audio_extension != ".wav":
        audio = AudioSegment.from_file(audio_path)
        wav_path = audio_path.replace(audio_extension, ".wav")
        audio.export(wav_path, format="wav")
        return wav_path
    return audio_path

# Function to extract text from an audio file
def extract_text_from_audio(file_path):
    try:
        # Convert to WAV format if not already in WAV
        wav_path = convert_to_wav(file_path)

        # Initialize recognizer
        recognizer = sr.Recognizer()

        # Use AudioFile to open and read the WAV file
        with sr.AudioFile(wav_path) as source:
            audio = recognizer.record(source)  # Record the entire file

        # Recognize speech using Google Speech Recognition
        text = recognizer.recognize_google(audio)
        return text

    except sr.UnknownValueError:
        return "Google Speech Recognition could not understand the audio."
    except sr.RequestError as e:
        return f"Could not request results from Google Speech Recognition service; {e}"
    except Exception as e:
        return f"Error processing audio file: {e}"
        


# Function to extract text from .py files
def extract_text_from_py(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            text = file.read()  # Read the entire Python file content
        return text
    except Exception as e:
        return f"Error reading Python file: {e}"
    
# Function to extract text from .wav (and other supported audio formats)
def extract_text_from_audio1(file_path):
    try:
        recognizer = sr.Recognizer()

        # Use AudioFile to open and read the audio file
        with sr.AudioFile(file_path) as source:
            audio = recognizer.record(source)  # Record the entire file

        # Recognize speech using Google Speech Recognition
        text = recognizer.recognize_google(audio)
        return text

    except sr.UnknownValueError:
        return "Google Speech Recognition could not understand the audio."
    except sr.RequestError as e:
        return f"Could not request results from Google Speech Recognition service; {e}"
    except Exception as e:
        return f"Error processing audio file: {e}"    
    
def extract_text_from_pptx(file_path):
    try:
        prs = Presentation(file_path)
        text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
        return text
    except Exception as e:
        return f"Error reading PPTX file: {e}"

def extract_text_from_zip(file_path):
    try:
        text = ""
        with tempfile.TemporaryDirectory() as temp_dir:
            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                zip_ref.extractall(temp_dir)
                for root, _, files in os.walk(temp_dir):
                    for file in files:
                        file_full_path = os.path.join(root, file)
                        file_ext = os.path.splitext(file_full_path)[1].lower()
                        if file_ext in extract_functions:
                            text += f"\n\nExtracted from {file}:\n\n" + extract_functions[file_ext](file_full_path)
                        else:
                            text += f"\n\nUnsupported file in zip: {file}"
        return text
    except Exception as e:
        return f"Error reading ZIP file: {e}"
      
    
# Function to extract text from .xml files
def extract_text_from_xml(file_path):
    try:
        tree = ET.parse(file_path)
        root = tree.getroot()
        text = ET.tostring(root, encoding='unicode', method='xml')
        return text
    except Exception as e:
        return f"Error reading XML file: {e}"

# Function to extract text from .xls files (using pandas)
def extract_text_from_xls(file_path):
    try:
        df = pd.read_excel(file_path, sheet_name=None)  # Read all sheets into a dictionary
        text = ""
        for sheet_name, sheet_df in df.items():
            text += f"Sheet: {sheet_name}\n"
            text += sheet_df.to_string(index=False) + "\n\n"  # Convert each sheet to a string
        return text
    except Exception as e:
        return f"Error reading XLS file: {e}. Make sure 'xlrd' is installed."    

def extract_text_from_pdb(file_path):
    with open(file_path, 'r') as file:
        pdb_text = file.read()
    return pdb_text


import tempfile
import os
from google.cloud import storage
import pandas as pd
import psycopg2
from psycopg2 import sql
from dotenv import load_dotenv

# Define your extraction functions here (e.g., extract_text_from_pdf, etc.)
# (Omitting function definitions for brevity, include them as defined earlier)

# Mapping extensions to their respective extraction functions
extract_functions = {
    ".pdf": extract_text_from_pdf,
    ".docx": extract_text_from_docx,
    ".xlsx": extract_text_from_excel,
    ".csv": extract_text_from_csv,
    ".json": extract_text_from_json,
    ".jsonld": extract_text_from_jsonld,
    ".jsonl": extract_text_from_jsonl,
    ".txt": extract_text_from_txt,
    ".jpg": extract_text_from_image,
    ".jpeg": extract_text_from_image,
    ".png": extract_text_from_image,
    ".py": extract_text_from_py,
    ".pptx": extract_text_from_pptx,
    ".mp3": extract_text_from_audio,
    ".wav": extract_text_from_audio1,
    ".zip": extract_text_from_zip,
    ".xls": extract_text_from_xls,
    ".xml": extract_text_from_xml,
    ".pdb": extract_text_from_pdb
}

# Load environment variables from .env file
load_dotenv()

# Load database configuration from environment variables
db_host = os.getenv("DB_HOST")
db_port = os.getenv("DB_PORT")
db_name = os.getenv("DB_NAME")
db_user = os.getenv("DB_USER")
db_password = os.getenv("DB_PASSWORD")

# GCP bucket name configuration
BUCKET_NAME = 'gaia_files'

# Function to connect to the Cloud SQL PostgreSQL instance using psycopg2
def connect_to_db():
    try:
        conn = psycopg2.connect(
            dbname=db_name,
            user=db_user,
            password=db_password,
            host=db_host,
            port=db_port
        )
        conn.autocommit = True  # Enable autocommit mode for the connection
        return conn
    except Exception as e:
        print(f"Error connecting to database: {e}")
        return None

# Function to add a new column to the table
def add_column_to_table(conn, table_name, column_name, column_type):
    try:
        # SQL query to add a new column if it does not exist
        add_column_query = sql.SQL("""
            ALTER TABLE {table}
            ADD COLUMN IF NOT EXISTS {column} {type};
        """).format(
            table=sql.Identifier(table_name),
            column=sql.Identifier(column_name),
            type=sql.SQL(column_type)
        )
        
        # Execute the query to add the column
        with conn.cursor() as cursor:
            cursor.execute(add_column_query)
        
        # Verify if the column was added successfully
        verify_column_query = sql.SQL("""
            SELECT column_name 
            FROM information_schema.columns 
            WHERE table_name = %s AND column_name = %s;
        """)
        
        with conn.cursor() as cursor:
            cursor.execute(verify_column_query, (table_name, column_name))
            result = cursor.fetchone()
        
        if result is None:
            print(f"Column '{column_name}' was not added to the table '{table_name}'.")
        else:
            print(f"Column '{column_name}' successfully added to the table '{table_name}'.")
    except Exception as e:
        print(f"Error adding column to table: {e}")

# Function to update the table with source text for each file
def update_table_with_source_text(conn, df, table_name):
    try:
        update_query = sql.SQL("""
            UPDATE {table}
            SET source_text = %s
            WHERE file_name = %s;
        """).format(table=sql.Identifier(table_name))
        
        with conn.cursor() as cursor:
            for _, row in df.iterrows():
                cursor.execute(update_query, (row['Extracted Text'], row['File_name']))
        print(f"Table '{table_name}' updated successfully with source text.")
    except Exception as e:
        print(f"Error updating table: {e}")

# Function to download all files from GCP bucket to a specific local directory
def download_files_to_directory(bucket_name, local_directory):
    try:
        # Initialize the Google Cloud Storage client
        client = storage.Client()
        bucket = client.bucket(bucket_name)
        
        # Create local directory if it doesn't exist
        if not os.path.exists(local_directory):
            os.makedirs(local_directory)
            print(f"Directory '{local_directory}' created.")
        
        # List all files in the bucket
        blobs = bucket.list_blobs()
        
        # Iterate through the files and download each one
        for blob in blobs:
            # Skip if it is a directory (e.g., if the name ends with a slash)
            if blob.name.endswith("/"):
                continue
            
            # Construct the local file path
            local_file_path = os.path.join(local_directory, os.path.basename(blob.name))
            
            # Download the file
            blob.download_to_filename(local_file_path)
            print(f"Downloaded {blob.name} to {local_file_path}")
    
    except Exception as e:
        print(f"Error downloading files from GCS: {e}")

# Function to extract text from a given directory of files (already defined by you)
def extract_text_from_directory(directory_path):
    # Files to ignore
    files_to_ignore = {"metadata.jsonl", "metadata.csv", ".DS_Store"}
    
    extracted_texts = {}
    if os.path.isdir(directory_path):
        for file_name in os.listdir(directory_path):
            # Skip files in the ignore list
            if file_name in files_to_ignore:
                continue
            
            file_path = os.path.join(directory_path, file_name)
            file_extension = os.path.splitext(file_path)[1].lower()
            if file_extension in extract_functions:
                extracted_texts[file_name] = extract_functions[file_extension](file_path)
            else:
                extracted_texts[file_name] = f"Unsupported file type: {file_extension}"
    else:
        return pd.DataFrame()  # Return an empty DataFrame if the directory is invalid

    # Convert results to a DataFrame
    return pd.DataFrame(list(extracted_texts.items()), columns=['File_name', 'Extracted Text'])

# Main workflow function
def main_workflow(bucket_name, table_name):
    # Define the local directory for downloading files
    local_dir = '/tmp/validation'  # Local path for downloaded files

    # Step 1: Download files from GCP bucket to the local directory
    download_files_to_directory(bucket_name, local_dir)

    # Step 2: Connect to the database
    conn = connect_to_db()
    if not conn:
        print("Failed to connect to the database.")
        return

    # Step 3: Add the 'source_text' column to the table
    add_column_to_table(conn, table_name, "source_text", "TEXT")

    # Step 4: Extract text from files in the local directory
    df_extracted_texts = extract_text_from_directory(local_dir)
    
    # Check if DataFrame is not empty
    if df_extracted_texts.empty:
        print("No valid files found for extraction.")
        return

    # Step 5: Update the table with the extracted text
    update_table_with_source_text(conn, df_extracted_texts, table_name)

    # Close the connection
    conn.close()

# Example usage
if __name__ == "__main__":
    table_name = 'validation'  # Replace with your actual table name
    main_workflow(BUCKET_NAME, table_name)