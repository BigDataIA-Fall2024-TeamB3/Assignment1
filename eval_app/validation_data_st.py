import streamlit as st
import os
import zipfile
import pandas as pd
from PyPDF2 import PdfReader
import docx
import json
import csv
from PIL import Image
import speech_recognition as sr
import pytesseract
import tempfile
from sqlalchemy import create_engine
from pptx import Presentation
from Bio import PDB
from dotenv import load_dotenv
import base64
import io
from pydub import AudioSegment
import os

# Load environment variables from .env file
load_dotenv()

# Load database configuration from environment variables
db_host = os.getenv("DB_HOST")
db_port = os.getenv("DB_PORT")
db_name = os.getenv("DB_NAME")
db_user = os.getenv("DB_USER")
db_password = os.getenv("DB_PASSWORD")

# Function to connect to the Cloud SQL PostgreSQL instance using SQLAlchemy
def connect_to_db():
    try:
        db_url = f'postgresql+psycopg2://{db_user}:{db_password}@{db_host}:{db_port}/{db_name}'
        engine = create_engine(db_url)
        return engine
    except Exception as e:
        st.error(f"Error connecting to database: {e}")
        return None

# Function to load data into a Pandas DataFrame
def load_data():
    engine = connect_to_db()
    if engine:
        query = "SELECT * FROM validation;"
        try:
            df = pd.read_sql(query, engine)
            return df
        except Exception as e:
            st.error(f"Error loading data: {e}")
    return None

# Define the extraction functions for all supported file types
def extract_text_from_pdf(file_path):
    try:
        with open(file_path, 'rb') as file:
            reader = PdfReader(file)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
        return text
    except Exception as e:
        return f"Error reading PDF file: {e}"

def extract_text_from_docx(file_path):
    try:
        doc = docx.Document(file_path)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        return text
    except Exception as e:
        return f"Error reading DOCX file: {e}"

def extract_text_from_excel(file_path):
    try:
        df = pd.read_excel(file_path, engine='openpyxl')
        text = df.to_string(index=False)
        return text
    except Exception as e:
        return f"Error reading Excel file: {e}. Make sure 'openpyxl' is installed."

def extract_text_from_csv(file_path):
    try:
        with open(file_path, 'r', newline='', encoding='utf-8') as file:
            reader = csv.reader(file)
            text = "\n".join([", ".join(row) for row in reader])
        return text
    except Exception as e:
        return f"Error reading CSV file: {e}"

def extract_text_from_json(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            data = json.load(file)
            text = json.dumps(data, indent=4)
        return text
    except Exception as e:
        return f"Error reading JSON file: {e}"
    
def extract_text_from_jsonld(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            data = json.load(file)
            text = json.dumps(data, indent=4)  # Format the JSON-LD data as a pretty-printed string
        return text
    except Exception as e:
        return f"Error reading JSON-LD file: {e}"    
    
# Function to extract text from .jsonl files
def extract_text_from_jsonl(file_path):
    try:
        text = ""
        with open(file_path, 'r', encoding='utf-8') as file:
            for line in file:
                try:
                    json_object = json.loads(line)
                    text += json.dumps(json_object, indent=4) + "\n"
                except json.JSONDecodeError as e:
                    text += f"Invalid JSON in line: {line}\nError: {e}\n"
        return text
    except Exception as e:
        return f"Error reading JSONL file: {e}"

def extract_text_from_txt(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            text = file.read()
        return text
    except Exception as e:
        return f"Error reading TXT file: {e}"

def extract_text_from_image(file_path):
    try:
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img)
        return text
    except Exception as e:
        return f"Error reading image file: {e}"


# Convert any audio file format to WAV using pydub
def convert_to_wav(audio_path):
    audio_extension = os.path.splitext(audio_path)[1].lower()
    if audio_extension != ".wav":
        audio = AudioSegment.from_file(audio_path)
        wav_path = audio_path.replace(audio_extension, ".wav")
        audio.export(wav_path, format="wav")
        return wav_path
    return audio_path

# Function to extract text from an audio file
def extract_text_from_audio(file_path):
    try:
        # Convert to WAV format if not already in WAV
        wav_path = convert_to_wav(file_path)

        # Initialize recognizer
        recognizer = sr.Recognizer()

        # Use AudioFile to open and read the WAV file
        with sr.AudioFile(wav_path) as source:
            audio = recognizer.record(source)  # Record the entire file

        # Recognize speech using Google Speech Recognition
        text = recognizer.recognize_google(audio)
        return text

    except sr.UnknownValueError:
        return "Google Speech Recognition could not understand the audio."
    except sr.RequestError as e:
        return f"Could not request results from Google Speech Recognition service; {e}"
    except Exception as e:
        return f"Error processing audio file: {e}"
        


# Function to extract text from .py files
def extract_text_from_py(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            text = file.read()  # Read the entire Python file content
        return text
    except Exception as e:
        return f"Error reading Python file: {e}"
    
# Function to extract text from .wav (and other supported audio formats)
def extract_text_from_audio1(file_path):
    try:
        recognizer = sr.Recognizer()

        # Use AudioFile to open and read the audio file
        with sr.AudioFile(file_path) as source:
            audio = recognizer.record(source)  # Record the entire file

        # Recognize speech using Google Speech Recognition
        text = recognizer.recognize_google(audio)
        return text

    except sr.UnknownValueError:
        return "Google Speech Recognition could not understand the audio."
    except sr.RequestError as e:
        return f"Could not request results from Google Speech Recognition service; {e}"
    except Exception as e:
        return f"Error processing audio file: {e}"    
    
def extract_text_from_pptx(file_path):
    try:
        prs = Presentation(file_path)
        text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
        return text
    except Exception as e:
        return f"Error reading PPTX file: {e}"

def extract_text_from_zip(file_path):
    try:
        text = ""
        with tempfile.TemporaryDirectory() as temp_dir:
            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                zip_ref.extractall(temp_dir)
                for root, _, files in os.walk(temp_dir):
                    for file in files:
                        file_full_path = os.path.join(root, file)
                        file_ext = os.path.splitext(file_full_path)[1].lower()
                        if file_ext in extract_functions:
                            text += f"\n\nExtracted from {file}:\n\n" + extract_functions[file_ext](file_full_path)
                        else:
                            text += f"\n\nUnsupported file in zip: {file}"
        return text
    except Exception as e:
        return f"Error reading ZIP file: {e}"
      
    
# Function to extract text from .xml files
def extract_text_from_xml(file_path):
    try:
        tree = ET.parse(file_path)
        root = tree.getroot()
        text = ET.tostring(root, encoding='unicode', method='xml')
        return text
    except Exception as e:
        return f"Error reading XML file: {e}"

# Function to extract text from .xls files (using pandas)
def extract_text_from_xls(file_path):
    try:
        df = pd.read_excel(file_path, sheet_name=None)  # Read all sheets into a dictionary
        text = ""
        for sheet_name, sheet_df in df.items():
            text += f"Sheet: {sheet_name}\n"
            text += sheet_df.to_string(index=False) + "\n\n"  # Convert each sheet to a string
        return text
    except Exception as e:
        return f"Error reading XLS file: {e}. Make sure 'xlrd' is installed."    

def extract_text_from_pdb(file_path):
    try:
        parser = PDB.PDBParser()
        structure = parser.get_structure('PDB_structure', file_path)
        text = f"PDB ID: {structure.id}\n"
        for model in structure:
            for chain in model:
                text += f"Chain: {chain.id}\n"
                for residue in chain:
                    text += f"Residue: {residue.resname} {residue.id}\n"
        return text
    except Exception as e:
        return f"Error reading PDB file: {e}"

# Mapping extensions to their respective extraction functions
extract_functions = {
    ".pdf": extract_text_from_pdf,
    ".docx": extract_text_from_docx,
    ".xlsx": extract_text_from_excel,
    ".csv": extract_text_from_csv,
    ".json": extract_text_from_json,
    ".jsonld": extract_text_from_jsonld,
    ".jsonl": extract_text_from_jsonl,
    ".txt": extract_text_from_txt,
    ".jpg": extract_text_from_image,
    ".jpeg": extract_text_from_image,
    ".png": extract_text_from_image,
    ".py": extract_text_from_py,
    ".pptx": extract_text_from_pptx,
    ".mp3": extract_text_from_audio,
    ".wav": extract_text_from_audio1,
    ".zip": extract_text_from_zip,
    ".xls": extract_text_from_xls,
    ".xml": extract_text_from_xml,
    ".pdb": extract_text_from_pdb
}

# Function to extract text from a given directory of files
def extract_text_from_directory(directory_path):
    extracted_texts = {}

    if os.path.isdir(directory_path):
        file_list = os.listdir(directory_path)
        for file_name in file_list:
            file_path = os.path.join(directory_path, file_name)
            file_extension = os.path.splitext(file_path)[1].lower()
            if file_extension in extract_functions:
                extracted_texts[file_name] = extract_functions[file_extension](file_path)
            else:
                extracted_texts[file_name] = f"Unsupported file type: {file_extension}"
    else:
        st.error("Invalid directory path.")
        return pd.DataFrame()
    

    # Convert results to a DataFrame
    extracted_texts_df = pd.DataFrame(list(extracted_texts.items()), columns=['File_name', 'Extracted Text'])
    return extracted_texts_df

# Function to display the Validation Data tab with file processing
def show():
    st.header("Validation Data")

    # Load the validation data from the database
    data = load_data()
    if data is not None:
        st.dataframe(data)
        st.subheader("Extract Text from Files")

        # Allow user to upload a directory
        directory_path = st.text_input("Enter the directory path to process files:", value="")
        if st.button("Process Files"):
            if directory_path and os.path.isdir(directory_path):
                extracted_texts_df = extract_text_from_directory(directory_path)
                st.write("Extracted Texts:")
                st.dataframe(extracted_texts_df)
                extracted_texts_df.to_csv("extracted_texts.csv", index=False)
                st.success("Text extraction completed and saved to 'extracted_texts.csv'.")
            else:
                st.error("Please enter a valid directory path.")
    else:
        st.write("No data to display.")

# Call the show function to display in the Streamlit app
if __name__ == "__main__":
    show()
